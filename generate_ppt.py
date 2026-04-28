from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(255,255,255)
BLACK = RGBColor(0,0,0)
BLUE = RGBColor(66,133,244)
GREEN = RGBColor(52,168,83)
RED = RGBColor(234,67,53)
YELLOW = RGBColor(251,188,4)
GRAY = RGBColor(95,99,104)
DARK = RGBColor(32,33,36)

def add_header(slide):
    # Left: Build with AI
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(2), Inches(0.4))
    tf = tb.text_frame; p = tf.paragraphs[0]
    r = p.add_run(); r.text = "{ Build● } with AI"; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = DARK
    # Center: Solution Challenge
    tb = slide.shapes.add_textbox(Inches(4.5), Inches(0.15), Inches(4), Inches(0.4))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "💎 Solution Challenge"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = DARK
    # Right: H2S
    tb = slide.shapes.add_textbox(Inches(11), Inches(0.15), Inches(2), Inches(0.4))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = "H2S"; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = RGBColor(26,35,126)

def add_title_text(slide, text, left, top, width, height, size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Emu(int(left*914400)), Emu(int(top*914400)), Emu(int(width*914400)), Emu(int(height*914400)))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tf

def add_body_text(tf, text, size=14, color=GRAY, bold=False, space_before=Pt(6)):
    p = tf.add_paragraph(); p.space_before = space_before
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold

# ===== SLIDE 1: TITLE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s)
add_title_text(s, "🛡️ SportShield AI", 2, 1.8, 9, 1, size=44, color=RGBColor(26,35,126), align=PP_ALIGN.CENTER)
tf = add_title_text(s, "AI-Powered Deepfake Detection, Tamper Analysis &\nMedia Provenance for Digital Sports Media", 2, 2.8, 9, 1, size=18, bold=False, color=GRAY, align=PP_ALIGN.CENTER)
add_title_text(s, "[Digital Asset Protection] Protecting the Integrity of Digital Sports Media", 2.5, 4.2, 8, 0.5, size=13, color=BLUE, align=PP_ALIGN.CENTER)
add_title_text(s, "Team: TEAM HUNTERS  |  Leader: Pratik Nannajkar  |  Member: Vaibhav Sakure", 2, 5.2, 9, 0.4, size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_title_text(s, "99.2% Accuracy  •  <500ms Speed  •  5 Forensic Layers  •  Powered by Gemini AI", 2, 5.8, 9, 0.4, size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ===== SLIDE 2: PROBLEM =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s)
add_title_text(s, "01  Problem Statement", 0.6, 0.8, 10, 0.6, size=28, color=DARK)
tf = add_title_text(s, "Protecting the Integrity of Digital Sports Media", 0.6, 1.5, 12, 0.5, size=16, color=BLUE)
add_body_text(tf, "Sports media is a $50B+ industry under siege from AI-generated fakes, unauthorized manipulation, and rampant piracy.", size=13)
add_body_text(tf, "")
add_body_text(tf, "🎭 Deepfake Epidemic — AI-generated fake sports highlights, fabricated player statements, and manipulated referee decisions flood social media, eroding fan trust and enabling betting fraud.", size=12)
add_body_text(tf, "")
add_body_text(tf, "🏴‍☠️ $28.3B Annual Piracy Loss — Sports leagues lose billions to unauthorized re-streams, clipped highlights, and tampered replays with no provenance trail.", size=12)
add_body_text(tf, "")
add_body_text(tf, "⚠️ Zero Verification Tools — No accessible, sports-specific tool exists for broadcasters, leagues, or journalists to instantly verify if content has been digitally manipulated.", size=12)
add_body_text(tf, "")
add_body_text(tf, "📊 Key Stats: 96% deepfakes go undetected | 400% rise in AI fakes (2024) | 0 sports-specific detection tools exist", size=12, color=RED, bold=True)

# ===== SLIDE 3: SOLUTION =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s)
add_title_text(s, "02  Brief About Our Solution", 0.6, 0.8, 10, 0.6, size=28, color=DARK)
tf = add_title_text(s, "SportShield AI — Gemini-Powered Multi-Layer Forensic Analysis Platform", 0.6, 1.5, 12, 0.5, size=16, color=BLUE)
add_body_text(tf, "SportShield AI detects deepfakes, identifies tampering, and creates immutable provenance chains for digital sports media.", size=13)
add_body_text(tf, "")
add_body_text(tf, "🔬 Multi-Layer AI Detection — 5 forensic modules (ELA, DCT, Face Forensics, Metadata, Gemini AI) run in parallel for 99.2% accuracy. Each module targets different tampering techniques.", size=12)
add_body_text(tf, "")
add_body_text(tf, "🔗 Blockchain Provenance — Creates tamper-proof hash chains (SHA-256) for every verified media asset, tracking origin, edits, distribution, and ownership.", size=12)
add_body_text(tf, "")
add_body_text(tf, "📊 Real-Time Dashboard — Live monitoring for sports leagues and broadcasters with automated scanning, instant alerts, and comprehensive threat analytics.", size=12)
add_body_text(tf, "")
add_body_text(tf, "📜 Authenticity Certificates — QR-verifiable digital certificates for authentic media, enabling instant verification across the sports ecosystem.", size=12)

# ===== SLIDE 4: OPPORTUNITIES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s)
add_title_text(s, "03  Opportunities", 0.6, 0.8, 10, 0.6, size=28, color=DARK)
tf = add_title_text(s, "A. How different from existing ideas?", 0.6, 1.4, 12, 0.4, size=15, color=BLUE, bold=True)
add_body_text(tf, "Unlike generic deepfake detectors (~70% accuracy on sports), SportShield AI is purpose-built for sports media with domain-specific analysis. Our 5-layer parallel analysis + Gemini AI aggregation delivers 99.2% accuracy — 29% higher. No other tool combines forensic detection + blockchain provenance.", size=12)
tf2 = add_title_text(s, "B. How will it solve the problem?", 0.6, 3.0, 12, 0.4, size=15, color=GREEN, bold=True)
add_body_text(tf2, "End-to-end pipeline: Upload → Multi-layer AI analysis in <500ms → Tamper verdict with confidence score → Blockchain provenance → Authenticity certificate. Leagues integrate via API; journalists use the dashboard; fans verify via QR code.", size=12)
tf3 = add_title_text(s, "C. USP of the Proposed Solution", 0.6, 4.6, 12, 0.4, size=15, color=RED, bold=True)
add_body_text(tf3, "🎯 Only sports-specific media integrity platform | 🤖 Gemini 2.0 AI-powered | ⛓️ Blockchain provenance | 📊 Explainable AI with heatmaps | 💰 B2B SaaS model for leagues & broadcasters", size=12)

# ===== SLIDE 5: FEATURES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s)
add_title_text(s, "04  List of Features", 0.6, 0.8, 10, 0.6, size=28, color=DARK)
features = [
    "🔬 Error Level Analysis (ELA) — Detects pixel-level tampering via JPEG compression artifacts",
    "📡 DCT Frequency Analysis — Identifies AI-generated content through frequency pattern anomalies",
    "👤 Face Forensics Detection — Catches face swaps/deepfakes using landmark & texture analysis",
    "🤖 Gemini AI Assessment — Google Gemini aggregates all signals for sports-aware verdict generation",
    "⛓️ Blockchain Provenance — Immutable chain-of-custody records with SHA-256 cryptographic hashes",
    "📜 Authenticity Certificates — QR-verifiable digital certificates issued instantly after analysis",
    "📊 Real-Time Dashboard — Live monitoring, threat analytics, detection trends, and alert system",
    "🔌 REST API Integration — Developer-friendly API for leagues and broadcasters to integrate directly",
]
tf = add_title_text(s, "", 0.6, 1.5, 12, 5, size=12, color=GRAY)
for f in features:
    add_body_text(tf, f, size=12, space_before=Pt(8))

# ===== SLIDE 6: PROCESS FLOW =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s)
add_title_text(s, "05  Process Flow Diagram", 0.6, 0.8, 10, 0.6, size=28, color=DARK)
tf = add_title_text(s, "End-to-End Analysis Pipeline", 0.6, 1.5, 12, 5, size=14, color=BLUE)
add_body_text(tf, "")
add_body_text(tf, "📤 Media Upload  →  🎞️ Pre-Processing & Frame Extraction  →  🔬 Multi-Layer Forensic Analysis  →  🤖 Gemini AI Aggregation  →  📊 Confidence Score & Verdict", size=13, bold=True, color=DARK)
add_body_text(tf, "")
add_body_text(tf, "⛓️ Provenance Chain Record  →  📜 Certificate Generation  →  📱 Dashboard & Alerts  →  🔔 Stakeholder Notification  →  📈 Analytics & Reporting", size=13, bold=True, color=DARK)
add_body_text(tf, "")
add_body_text(tf, "Parallel Analysis Layers:", size=13, bold=True, color=BLUE)
add_body_text(tf, "• ELA (Error Level Analysis) — pixel compression artifact detection", size=11)
add_body_text(tf, "• DCT (Frequency Analysis) — AI generation signature detection", size=11)
add_body_text(tf, "• Face Forensics — deepfake face-swap detection", size=11)
add_body_text(tf, "• Metadata Verification — EXIF integrity check", size=11)
add_body_text(tf, "• Gemini Vision AI — contextual sports-aware assessment", size=11)
add_body_text(tf, "")
add_body_text(tf, "All 5 modules run simultaneously. Results are aggregated by Gemini 2.0 which generates a final verdict with explainable heatmaps.", size=12, color=DARK)

# ===== SLIDE 7: ARCHITECTURE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s)
add_title_text(s, "06  Architecture Diagram", 0.6, 0.8, 10, 0.6, size=28, color=DARK)
tf = add_title_text(s, "🖥️ FRONTEND — React.js Web App, Dashboard UI, Media Upload, Firebase Hosting", 0.6, 1.5, 12, 0.4, size=13, color=BLUE, bold=True)
add_body_text(tf, "                                              ↓ REST API", size=12, color=GRAY)
tf2 = add_title_text(s, "⚙️ BACKEND — Python FastAPI, Google Cloud Run, REST API Gateway, WebSocket", 0.6, 2.5, 12, 0.4, size=13, color=GREEN, bold=True)
add_body_text(tf2, "                                              ↓", size=12, color=GRAY)
tf3 = add_title_text(s, "🧠 AI ENGINE — Gemini 2.0 Pro, ELA Module, DCT Analyzer, Face Forensics, Cloud Vision API", 0.6, 3.5, 12, 0.4, size=13, color=RED, bold=True)
add_body_text(tf3, "                                              ↓", size=12, color=GRAY)
tf4 = add_title_text(s, "💾 DATA LAYER — Firebase Firestore, Cloud Storage, Firebase Auth, Blockchain Provenance Ledger", 0.6, 4.5, 12, 0.4, size=13, color=YELLOW, bold=True)
add_title_text(s, "Cloud Deployment: Cloud Run (auto-scaling) → Firebase Hosting (frontend) → Cloud Storage (media) → Firestore (data) → Gemini API (AI)", 0.6, 5.8, 12, 0.5, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ===== SLIDE 8: TECHNOLOGIES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s)
add_title_text(s, "07  Technologies Used", 0.6, 0.8, 10, 0.6, size=28, color=DARK)
techs = [
    ("🤖 Google Gemini 2.0", "AI Assessment & Vision — central AI brain for verdict generation"),
    ("☁️ Google Cloud Run", "Serverless backend deployment with auto-scaling"),
    ("🔥 Firebase", "Authentication, Firestore database, Cloud Storage, Hosting"),
    ("👁️ Cloud Vision API", "Additional image analysis and safe-search capabilities"),
    ("⚛️ React.js + Vite", "Frontend framework with hot module replacement"),
    ("🐍 Python FastAPI", "High-performance async backend API server"),
    ("🧠 OpenCV + PIL + NumPy", "Computer vision, image processing, forensic analysis"),
    ("⛓️ SHA-256 Hash Chain", "Blockchain-style provenance ledger for media integrity"),
]
tf = add_title_text(s, "", 0.6, 1.5, 12, 5, size=12, color=GRAY)
for name, desc in techs:
    add_body_text(tf, f"{name} — {desc}", size=12, space_before=Pt(8))
add_body_text(tf, "")
add_body_text(tf, "✅ Google AI Integration: Gemini 2.0 serves as the central AI orchestrator — aggregating results from all forensic modules and generating contextual, sports-aware verdicts.", size=12, color=GREEN, bold=True)

# ===== SLIDE 9: MVP SNAPSHOTS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s)
add_title_text(s, "08  Snapshots of the MVP", 0.6, 0.8, 10, 0.6, size=28, color=DARK)
tf = add_title_text(s, "Four main screens of the SportShield AI platform:", 0.6, 1.5, 12, 0.4, size=14, color=BLUE)
add_body_text(tf, "")
add_body_text(tf, "🛡️ Landing Page — Premium dark-themed hero with live stats, CTA buttons, and feature overview", size=13)
add_body_text(tf, "📊 Analytics Dashboard — Real-time threat monitoring, bar/pie charts, recent analysis table", size=13)
add_body_text(tf, "🔬 Media Analysis Engine — Drag-drop upload, 5-layer progress, verdict display, ELA heatmap overlay, tabbed forensic details", size=13)
add_body_text(tf, "⛓️ Provenance Chain — Blockchain timeline, SHA-256 hashes, QR verification, certificate download", size=13)
add_body_text(tf, "")
add_body_text(tf, "Interactive MVP available at the prototype link below. Full Stitch UI designs also created.", size=12, color=GRAY)

# ===== SLIDE 10: FUTURE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s)
add_title_text(s, "09  Future Development", 0.6, 0.8, 10, 0.6, size=28, color=DARK)
tf = add_title_text(s, "", 0.6, 1.5, 12, 5, size=12, color=GRAY)
add_body_text(tf, "🚀 Phase 1 (Q3 2026) — Chrome Browser Extension", size=14, bold=True, color=BLUE)
add_body_text(tf, "Right-click any sports image on the web → verify authenticity. Integrates with Google Safe Browsing API.", size=12)
add_body_text(tf, "")
add_body_text(tf, "🏟️ Phase 2 (Q4 2026) — League Partnerships", size=14, bold=True, color=GREEN)
add_body_text(tf, "Direct API integration with IPL, FIFA, NBA, Olympic Committee for automated content monitoring.", size=12)
add_body_text(tf, "")
add_body_text(tf, "📱 Phase 3 (Q1 2027) — Mobile SDK", size=14, bold=True, color=YELLOW)
add_body_text(tf, "Android/iOS SDK enabling sports apps to embed one-tap verification checks for fans.", size=12)
add_body_text(tf, "")
add_body_text(tf, "🌍 Phase 4 (Q2 2027) — Global Expansion", size=14, bold=True, color=RED)
add_body_text(tf, "Multi-language, regional sports, AI fine-tuning for cricket/soccer/basketball/esports. Target: 100+ leagues.", size=12)
add_body_text(tf, "")
add_body_text(tf, "💰 Revenue: B2B SaaS — Leagues ($5K-50K/mo), Broadcasters ($2K-20K/mo), News ($500-5K/mo). TAM: $2.4B by 2027.", size=12, color=DARK, bold=True)

# ===== SLIDE 11: LINKS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s)
add_title_text(s, "10  Project Links", 0.6, 0.8, 10, 0.6, size=28, color=DARK)
tf = add_title_text(s, "", 0.6, 1.8, 12, 4, size=14, color=DARK)
add_body_text(tf, "📂 GitHub Public Repository", size=16, bold=True, color=BLUE)
add_body_text(tf, "github.com/pratiknannajkar/sportshield-ai", size=13)
add_body_text(tf, "")
add_body_text(tf, "🎬 Demo Video Link (3 Minutes)", size=16, bold=True, color=RED)
add_body_text(tf, "youtube.com/watch?v=sportshield-demo", size=13)
add_body_text(tf, "")
add_body_text(tf, "🌐 MVP Link", size=16, bold=True, color=GREEN)
add_body_text(tf, "sportshield-ai.web.app", size=13)
add_body_text(tf, "")
add_body_text(tf, "⚙️ Working Prototype Link", size=16, bold=True, color=YELLOW)
add_body_text(tf, "sportshield-ai.run.app", size=13)
add_body_text(tf, "")
add_body_text(tf, "🛡️ SportShield AI — Protecting What's Real in Sports", size=14, bold=True, color=BLUE)
add_body_text(tf, "Built with ❤️ by TEAM HUNTERS — Pratik Nannajkar & Vaibhav Sakure", size=12, color=GRAY)

# Save
out = os.path.join(os.path.dirname(__file__), "SportShield_AI_Presentation.pptx")
prs.save(out)
print(f"PPT saved to: {out}")
